#!/usr/bin/env python3
"""Regression tests added after the full audit.

They lock the fixes that closed the audit findings: the OOXML core-date pinning that
makes the xlsx byte-deterministic (verified without an Office backend, by exercising
the port directly), the HTML escaping of hostile labels, and the input-validation
error paths that several modules gained (missing run timestamp, an STL period below a
cycle, missing measured clicks, an unknown module, malformed brand input).
"""

from __future__ import annotations

import io
import sys
import zipfile
from datetime import datetime, timezone
from pathlib import Path

import pytest

PLUGIN_ROOT = Path(__file__).resolve().parent.parent
if str(PLUGIN_ROOT) not in sys.path:
    sys.path.insert(0, str(PLUGIN_ROOT))

SAMPLES = PLUGIN_ROOT / "assets" / "samples"
FIXED_TIMESTAMP = "2026-06-08T00:00:00+00:00"


# --- determinism: the OOXML core-date pinning (backend-free) -----------------


def _zip_with_core(modified: str, created: str) -> bytes:
    core = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<cp:coreProperties xmlns:cp="x" xmlns:dcterms="y" xmlns:xsi="z">'
        f'<dcterms:created xsi:type="dcterms:W3CDTF">{created}</dcterms:created>'
        f'<dcterms:modified xsi:type="dcterms:W3CDTF">{modified}</dcterms:modified>'
        "</cp:coreProperties>"
    )
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w") as z:
        z.writestr("docProps/core.xml", core)
        z.writestr("[Content_Types].xml", "<Types/>")
    return buffer.getvalue()


def test_normalize_zip_pins_core_dates(tmp_path: Path) -> None:
    # openpyxl overwrites the modified date with the wall clock at save; the port must
    # rewrite it to the pinned run timestamp so two renders never differ on the date.
    from querymantic.ports import ooxml

    path = tmp_path / "doc.xlsx"
    path.write_bytes(_zip_with_core("2031-09-09T12:34:56Z", "2031-09-09T12:34:56Z"))
    ts = datetime(2026, 6, 8, 0, 0, 0, tzinfo=timezone.utc)
    ooxml.normalize_zip(path, core_timestamp=ts)
    core = zipfile.ZipFile(io.BytesIO(path.read_bytes())).read("docProps/core.xml")
    text = core.decode("utf-8")
    assert "2026-06-08T00:00:00Z" in text
    assert "2031-09-09T12:34:56Z" not in text


def test_parse_timestamp_accepts_zulu_suffix() -> None:
    # The pinned timestamp is stored with a trailing 'Z'; before Python 3.11
    # datetime.fromisoformat rejects it, which used to fall back silently to 1980.
    # The parser must normalise the 'Z' so every supported Python yields the same date.
    from modules.output_forge import _parse_timestamp

    parsed = _parse_timestamp("2026-06-08T00:00:00Z")
    assert parsed.year == 2026 and parsed.month == 6 and parsed.day == 8
    # The offset form parses to the same instant.
    assert _parse_timestamp("2026-06-08T00:00:00+00:00") == parsed


def test_parse_timestamp_rejects_unparsable_value() -> None:
    # A non-empty but malformed pinned timestamp must raise, never be replaced by an
    # invented date that would silently stamp the wrong year on every document.
    from modules.output_forge import OutputForgeError, _parse_timestamp

    with pytest.raises(OutputForgeError):
        _parse_timestamp("not-a-timestamp")
    with pytest.raises(OutputForgeError):
        _parse_timestamp("")


def test_xlsx_core_dates_pinned_not_1980(tmp_path: Path) -> None:
    # End-to-end guard for A1: with a 'Z'-suffixed pinned timestamp, the rendered
    # workbook's core.xml dates must be the pinned 2026 date, not the 1980 fallback.
    # Needs the openpyxl backend, so it is exercised in CI (and skipped where absent).
    openpyxl = pytest.importorskip("openpyxl")
    import zipfile

    from modules.output_forge import _parse_timestamp
    from modules.output_forge.brand import resolve_brand
    from modules.output_forge.dashboard_xlsx import render_workbook
    from querymantic.ports import ooxml

    view = {
        "header": {"label": "Audit"},
        "corpus": {
            "total_keywords": 1,
            "aio_eligibility_share": 0.1,
            "geo_opportunity_share": 0.1,
            "demand_opportunity_score": 1.0,
            "intent_split": [],
        },
        "clusters": [],
        "winnable": None,
        "gaps": [],
    }
    ts = _parse_timestamp("2026-06-08T00:00:00Z")
    out = tmp_path / "dash.xlsx"
    render_workbook(view, resolve_brand(None), out, ts)
    ooxml.normalize_zip(out, core_timestamp=ts)
    core = zipfile.ZipFile(out).read("docProps/core.xml").decode("utf-8")
    assert "2026-06-08T00:00:00Z" in core
    assert "1980-01-01" not in core
    assert openpyxl.__version__  # backend really exercised


def test_pin_core_dates_is_prefix_agnostic() -> None:
    # B1: the date pinning must not be bound to the literal 'dcterms:' prefix. A
    # different prefix (or none) used to leave the wall-clock date in place silently.
    from querymantic.ports import ooxml

    ts = datetime(2026, 6, 8, 0, 0, 0, tzinfo=timezone.utc)
    cases = (
        '<r xmlns:dct="y"><dct:created>2031-09-09T12:34:56Z</dct:created>'
        "<dct:modified>2031-09-09T12:34:56Z</dct:modified></r>",
        "<r><created>2031-01-01T00:00:00Z</created>"
        "<modified>2031-01-01T00:00:00Z</modified></r>",
        '<r xmlns:dcterms="y"><dcterms:modified>2031-01-01T00:00:00Z</dcterms:modified>'
        "<dcterms:created>2031-01-01T00:00:00Z</dcterms:created></r>",
    )
    for core in cases:
        out = ooxml._pin_core_dates(core.encode("utf-8"), ts).decode("utf-8")
        assert "2026-06-08T00:00:00Z" in out
        assert "2031" not in out


def test_pin_core_dates_raises_when_nothing_to_pin() -> None:
    # B1: if the core part carries no date element, pinning must raise rather than
    # silently no-op and let a wall-clock date survive.
    from querymantic.ports import ooxml

    ts = datetime(2026, 6, 8, 0, 0, 0, tzinfo=timezone.utc)
    with pytest.raises(ooxml.OoxmlError):
        ooxml._pin_core_dates(b"<r><title>no dates here</title></r>", ts)


def _state_for_forge(tmp_path: Path) -> dict:
    from querymantic import pipeline

    return pipeline.run_pipeline(
        PLUGIN_ROOT,
        [SAMPLES],
        tmp_path / "run.json",
        modules_to_run=(
            "entity_web",
            "fan_out_radar",
            "citation_grid",
            "click_ceiling",
        ),
        generated_at=FIXED_TIMESTAMP,
    )


# --- P1/P2: deliverable polish -------------------------------------------------


def test_displayed_scores_have_at_most_one_decimal(tmp_path: Path) -> None:
    # P1: the readiness scores must not carry false 4-decimal precision; they show one
    # decimal, coherent with the dashboard cards.
    from modules.output_forge.model import build_view
    from querymantic import pipeline

    state = pipeline.run_pipeline(
        PLUGIN_ROOT,
        [SAMPLES],
        tmp_path / "run.json",
        modules_to_run=(
            "entity_web",
            "fan_out_radar",
            "citation_grid",
            "click_ceiling",
        ),
        generated_at=FIXED_TIMESTAMP,
    )
    view = build_view(state)
    mean = view["readiness"]["mean_expected_readiness"]
    assert mean == round(mean, 1)
    assert view["corpus"]["demand_opportunity_score"] == round(
        view["corpus"]["demand_opportunity_score"], 1
    )
    for c in view["top_clusters"]:
        if "expected_readiness" in c:
            assert c["expected_readiness"] == round(c["expected_readiness"], 1)
        if "expected_share" in c:
            assert c["expected_share"] == round(c["expected_share"], 1)


def test_band_chart_value_label_is_not_clipped() -> None:
    # P2: a wide "low to high" band label must fit inside the chart, not run off the
    # right edge as it did when the value column was too narrow.
    import re

    from modules.output_forge import dashboard_html as dh

    colors = {
        "text": "#000000",
        "band": "#cccccc",
        "accent": "#ff0000",
        "muted": "#888888",
        "secondary": "#0000ff",
    }
    svg = dh._band_chart([("organic", 318122, 1820400)], colors)
    width = int(re.search(r'viewBox="0 0 (\d+)', svg).group(1))
    value_x = int(re.findall(r'<text x="(\d+)"[^>]*font-size="12"', svg)[-1])
    label = "318,122 to 1,820,400"
    # 7 user units per character is a conservative upper bound at font-size 12.
    assert value_x + len(label) * 7 <= width, "band label would be clipped"


# --- B4: scoring parameters are surfaced and overridable -----------------------


def test_geo_nudge_is_applied_and_overridable() -> None:
    from modules.citation_grid import _eligibility_value

    kw = {
        "scopes": {
            "aio_eligibility": {"score": 50},
            "geo_opportunity": {"label": "dual"},
        }
    }
    # Default nudge lifts a GEO-routed query by 0.1; an override changes it.
    assert _eligibility_value([0], [kw], 0.0) == 0.5
    assert _eligibility_value([0], [kw], 0.4) == round(0.9, 6)


def test_freshness_states_are_overridable() -> None:
    from modules.citation_grid import _freshness_value

    assert _freshness_value(0, {0: "rising"}) == 1.0
    assert _freshness_value(0, {0: "rising"}, {"rising": 0.2}) == 0.2


def test_citation_grid_surfaces_scoring_params(tmp_path: Path) -> None:
    from modules.citation_grid import citation_grid
    from querymantic import pipeline

    state = pipeline.run_pipeline(
        PLUGIN_ROOT,
        [SAMPLES],
        tmp_path / "run.json",
        modules_to_run=("entity_web", "fan_out_radar"),
        generated_at=FIXED_TIMESTAMP,
    )
    citation_grid(state, params={"geo_eligibility_nudge": 0.25})
    params = state["modules"]["citation_grid"]["params"]
    assert params["geo_eligibility_nudge"] == 0.25
    assert "fresh_states" in params and "provenance" in params


def test_click_ceiling_surfaces_thresholds_and_override_bites(tmp_path: Path) -> None:
    from modules.click_ceiling import click_ceiling
    from querymantic import pipeline

    def _run(params):
        state = pipeline.run_pipeline(
            PLUGIN_ROOT,
            [SAMPLES],
            tmp_path / "run.json",
            modules_to_run=("entity_web", "fan_out_radar", "citation_grid"),
            generated_at=FIXED_TIMESTAMP,
        )
        click_ceiling(state, params=params)
        return state["modules"]["click_ceiling"]

    # Max trim (every covered cluster counts AI-heavy and low-coverage) vs no trim.
    trimmed = _run({"ai_heavy_share": 0.0, "low_coverage_tau": 1.0})
    untrimmed = _run({"ai_heavy_share": 1.0, "low_coverage_tau": 0.0})
    p = trimmed["params"]
    assert p["ai_heavy_share"] == 0.0 and p["low_coverage_tau"] == 1.0
    assert "provenance" in p
    # The coverage trim only ever lowers the upper endpoint, so max-trim <= no-trim.
    assert (
        trimmed["summary"]["winnable_band"][1]
        < untrimmed["summary"]["winnable_band"][1]
    )


def _output_forge_module():
    # The package defines a function named output_forge that shadows the submodule in
    # an `import ... as` binding, so reach the module object explicitly to monkeypatch it.
    import importlib

    return importlib.import_module("modules.output_forge")


def test_missing_backend_is_a_clean_skip(tmp_path: Path, monkeypatch) -> None:
    # B3: an ImportError from the render path is the missing-backend case and must be a
    # recorded skip, not a failure, so the suite still degrades to HTML.
    from querymantic.ports import ooxml

    of = _output_forge_module()
    monkeypatch.setattr(
        ooxml,
        "ooxml_capabilities",
        lambda: {"pptx": False, "docx": False, "xlsx": True},
    )

    def _raise_import(*a, **k):
        raise ModuleNotFoundError("no openpyxl")

    monkeypatch.setattr(of, "_render_ooxml", _raise_import)
    state = _state_for_forge(tmp_path)
    of.output_forge(state, out_dir=tmp_path / "forge", formats=("html", "xlsx"))
    forge = state["modules"]["output_forge"]
    assert "html" in {a["format"] for a in forge["artifacts"]}
    xlsx_skip = next(s for s in forge["skipped"] if s["format"] == "xlsx")
    assert "not importable" in xlsx_skip["reason"]


def test_rendering_error_is_not_hidden_as_skip(tmp_path: Path, monkeypatch) -> None:
    # B3: a genuine rendering error (not a missing backend) must surface as an error,
    # never be downgraded to a silent skip with an obscure reason.
    from querymantic.ports import ooxml

    of = _output_forge_module()
    monkeypatch.setattr(
        ooxml,
        "ooxml_capabilities",
        lambda: {"pptx": False, "docx": False, "xlsx": True},
    )

    def _raise_value(*a, **k):
        raise ValueError("a real rendering bug")

    monkeypatch.setattr(of, "_render_ooxml", _raise_value)
    state = _state_for_forge(tmp_path)
    with pytest.raises(of.OutputForgeError):
        of.output_forge(state, out_dir=tmp_path / "forge", formats=("html", "xlsx"))


def test_short_hex_expands_to_six_digits() -> None:
    # B2: a 3-digit hex (#abc) used to pass validation then break the Office renderers.
    # resolve_brand now expands it once, and office_hex returns the 6-digit form.
    from modules.output_forge.brand import office_hex, resolve_brand

    brand = resolve_brand({"colors": {"primary": "#abc"}})
    assert brand["colors"]["primary"] == "#aabbcc"
    assert office_hex(brand["colors"]["primary"]) == "AABBCC"
    # office_hex is robust to a bare short value from any caller.
    assert office_hex("#abc") == "AABBCC"
    assert office_hex("#1F3A5F") == "1F3A5F"


def test_short_hex_renders_in_every_format(tmp_path: Path) -> None:
    # B2 end-to-end: with a 3-digit brand colour, all renderers must succeed and agree
    # on the expanded colour, instead of html rendering while xlsx/docx raise. The
    # Office formats need their backends, so they run in CI / on 3.10 and skip elsewhere.
    from modules.output_forge.brand import resolve_brand
    from modules.output_forge.dashboard_html import render_html
    from modules.output_forge.model import build_view
    from querymantic import pipeline

    state = pipeline.run_pipeline(
        PLUGIN_ROOT,
        [SAMPLES],
        tmp_path / "run.json",
        modules_to_run=(
            "entity_web",
            "fan_out_radar",
            "citation_grid",
            "click_ceiling",
        ),
        generated_at=FIXED_TIMESTAMP,
    )
    view = build_view(state)
    brand = resolve_brand({"colors": {"primary": "#abc"}})
    # HTML is backend-free and must render with the short form expanded.
    html = render_html(view, brand).decode("utf-8")
    assert "#aabbcc" in html

    from querymantic.ports import ooxml

    caps = ooxml.ooxml_capabilities()
    if not (caps["xlsx"] and caps["docx"]):
        pytest.skip("Office backends not installed")
    from datetime import datetime, timezone

    from modules.output_forge.audit_docx import render_audit
    from modules.output_forge.dashboard_xlsx import render_workbook

    ts = datetime(2026, 6, 8, tzinfo=timezone.utc)
    render_workbook(view, brand, tmp_path / "d.xlsx", ts)
    render_audit(view, brand, tmp_path / "d.docx", ts)
    assert (tmp_path / "d.xlsx").is_file() and (tmp_path / "d.docx").is_file()


def test_safe_cell_neutralises_formula_leaders() -> None:
    # Backend-free unit guard for M1: a string starting with a formula leader is
    # prefixed so Excel reads it as text; a clean string and a number pass through.
    from modules.output_forge.dashboard_xlsx import _safe_cell

    for payload in ("=1+2", "=cmd|'/c calc'!A1", "+1", "-1", "@SUM(1)", "\t=1"):
        assert _safe_cell(payload) == "'" + payload
    assert _safe_cell("running shoes") == "running shoes"
    assert _safe_cell(168000) == 168000
    assert _safe_cell(None) is None


def test_xlsx_has_no_active_formula_from_input(tmp_path: Path) -> None:
    # End-to-end guard for M1: hostile cluster heads / entities / title must never
    # render as a live formula in the delivered workbook. Needs openpyxl (CI / 3.10).
    openpyxl = pytest.importorskip("openpyxl")
    from datetime import datetime, timezone

    from modules.output_forge.brand import resolve_brand
    from modules.output_forge.dashboard_xlsx import render_workbook

    payloads = ["=1+2", "=cmd|'/c calc'!A1", '=HYPERLINK("http://evil")']
    view = {
        "header": {"label": payloads[2]},
        "corpus": {
            "total_keywords": 1,
            "aio_eligibility_share": 0.1,
            "geo_opportunity_share": 0.1,
            "demand_opportunity_score": 1.0,
            "intent_split": [{"intent": "=BAD", "count": 1}],
        },
        "clusters": [
            {
                "head": payloads[0],
                "size": 1,
                "volume_total": 10,
                "dominant_intent": "x",
                "winnable_band": [1, 2],
            }
        ],
        "winnable": None,
        "gaps": [
            {
                "entity": payloads[1],
                "demand_volume": 5,
                "suggested_cluster_head": "=EVIL",
            }
        ],
    }
    out = tmp_path / "dash.xlsx"
    render_workbook(
        view, resolve_brand(None), out, datetime(2026, 6, 8, tzinfo=timezone.utc)
    )
    wb = openpyxl.load_workbook(out)
    active = [
        (ws.title, cell.coordinate, cell.value)
        for ws in wb.worksheets
        for row in ws.iter_rows()
        for cell in row
        if cell.data_type == "f"
    ]
    assert active == [], f"active formulas leaked into the workbook: {active}"


def test_pptx_is_deterministic_and_safe_with_hostile_text(tmp_path: Path) -> None:
    # The pptx format was the one deliverable without an explicit determinism /
    # text-injection test. With python-pptx in CI: two renders are byte-identical
    # after archive normalisation, and a hostile cluster head is stored as literal
    # text in a valid file (python-pptx escapes it; pptx has no formula concept).
    pytest.importorskip("pptx")
    from datetime import datetime, timezone

    from pptx import Presentation

    from modules.output_forge.brand import resolve_brand
    from modules.output_forge.deck import render_deck
    from modules.output_forge.model import build_view
    from querymantic import pipeline
    from querymantic.ports import ooxml

    state = pipeline.run_pipeline(
        PLUGIN_ROOT,
        [SAMPLES],
        tmp_path / "run.json",
        modules_to_run=(
            "entity_web",
            "fan_out_radar",
            "citation_grid",
            "click_ceiling",
        ),
        generated_at=FIXED_TIMESTAMP,
    )
    view = build_view(state)
    hostile = "<script>alert(1)</script> =1+2 ]]>&"
    if view["top_clusters"]:
        view["top_clusters"][0]["head"] = hostile
    brand = resolve_brand(None)
    ts = datetime(2026, 6, 8, tzinfo=timezone.utc)

    a, b = tmp_path / "a.pptx", tmp_path / "b.pptx"
    for path in (a, b):
        render_deck(view, brand, path, ts)
        ooxml.normalize_zip(path, core_timestamp=ts)
    assert a.read_bytes() == b.read_bytes(), "pptx is not byte-deterministic"

    # The file is valid and the hostile text is present as literal text (escaped by
    # python-pptx into the XML, not breaking it).
    prs = Presentation(str(a))
    seen = any(
        hostile in shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert seen, "hostile cluster head not found as literal text in the deck"


def test_normalize_zip_without_timestamp_leaves_dates(tmp_path: Path) -> None:
    # When no timestamp is passed the port only normalises member metadata, not the
    # core dates, so the call stays backward-compatible.
    from querymantic.ports import ooxml

    path = tmp_path / "doc.xlsx"
    path.write_bytes(_zip_with_core("2031-09-09T12:34:56Z", "2031-09-09T12:34:56Z"))
    ooxml.normalize_zip(path)
    text = zipfile.ZipFile(io.BytesIO(path.read_bytes())).read("docProps/core.xml")
    assert b"2031-09-09T12:34:56Z" in text


# --- HTML escaping of a hostile label ----------------------------------------


def test_html_dashboard_escapes_hostile_label(tmp_path: Path) -> None:
    from modules.output_forge.brand import resolve_brand
    from modules.output_forge.dashboard_html import render_html
    from modules.output_forge.model import build_view
    from querymantic import pipeline

    state = pipeline.run_pipeline(
        PLUGIN_ROOT,
        [SAMPLES],
        tmp_path / "run.json",
        modules_to_run=(
            "entity_web",
            "fan_out_radar",
            "citation_grid",
            "click_ceiling",
        ),
        generated_at=FIXED_TIMESTAMP,
    )
    view = build_view(state)
    hostile = "<script>alert(1)</script>"
    view["corpus"]["intent_split"][0]["intent"] = hostile
    if view["top_clusters"]:
        view["top_clusters"][0]["head"] = hostile
    html = render_html(view, resolve_brand(None)).decode("utf-8")
    assert "<script>alert(1)</script>" not in html
    assert "&lt;script&gt;" in html


# --- input-validation error paths --------------------------------------------


def test_demand_pulse_rejects_stl_period_below_two() -> None:
    from modules.demand_pulse import ModuleError, _config

    with pytest.raises(ModuleError):
        _config({"stl_period": 1})


def test_demand_pulse_rejects_zero_momentum_window() -> None:
    from modules.demand_pulse import ModuleError, _config

    with pytest.raises(ModuleError):
        _config({"momentum_window": 0})


def test_live_wire_rejects_missing_clicks() -> None:
    from modules.live_wire import LiveWireError, _parse_search_console

    block = {"rows": [{"query": "running shoes", "impressions": 100}]}
    with pytest.raises(LiveWireError):
        _parse_search_console(block, Path("capture.json"))


def test_live_wire_keeps_a_real_zero() -> None:
    # A measured zero must survive, unlike a missing field, which raises above.
    from modules.live_wire import _parse_search_console

    block = {"rows": [{"query": "x", "clicks": 0, "impressions": 0}]}
    parsed = _parse_search_console(block, Path("capture.json"))
    assert parsed["rows"][0]["clicks"] == 0.0


def test_fan_out_requires_generated_at() -> None:
    from modules.fan_out_radar import ModuleError, _reference_year

    with pytest.raises(ModuleError):
        _reference_year({"querymantic": {"generated_at": ""}})


def test_reference_year_reads_the_stamp() -> None:
    from modules.fan_out_radar import _reference_year

    assert _reference_year({"querymantic": {"generated_at": FIXED_TIMESTAMP}}) == 2026


def test_pipeline_rejects_unknown_module(tmp_path: Path) -> None:
    from querymantic.pipeline import PipelineError, run_pipeline

    with pytest.raises(PipelineError):
        run_pipeline(
            PLUGIN_ROOT,
            [SAMPLES],
            tmp_path / "run.json",
            modules_to_run=("no_such_module",),
            generated_at=FIXED_TIMESTAMP,
        )


def test_load_brand_rejects_bad_json(tmp_path: Path) -> None:
    from modules.output_forge.brand import BrandError, load_brand

    bad = tmp_path / "brand.json"
    bad.write_text("{ not valid json", encoding="utf-8")
    with pytest.raises(BrandError):
        load_brand(bad)


def test_brand_rejects_bad_hex() -> None:
    from modules.output_forge.brand import BrandError, resolve_brand

    with pytest.raises(BrandError):
        resolve_brand({"name": "Acme", "colors": {"primary": '"/><script>'}})


def test_brand_resolution_is_idempotent() -> None:
    from modules.output_forge.brand import resolve_brand

    once = resolve_brand({"name": "Acme", "colors": {"primary": "#ABCDEF"}})
    twice = resolve_brand(once)
    assert once == twice


# --- loader error paths ------------------------------------------------------


def test_load_series_rejects_missing_file(tmp_path: Path) -> None:
    from modules.demand_pulse import DemandPulseError, load_series

    with pytest.raises(DemandPulseError):
        load_series(tmp_path / "no_such_series.csv")


def test_load_ctr_table_rejects_missing_file(tmp_path: Path) -> None:
    from modules.click_ceiling import ModuleError, load_ctr_table

    with pytest.raises(ModuleError):
        load_ctr_table(tmp_path / "no_such_table.json")


def test_load_gazetteer_rejects_missing_file(tmp_path: Path) -> None:
    from modules.language_layer import LanguageLayerError, load_gazetteer

    with pytest.raises(LanguageLayerError):
        load_gazetteer(tmp_path / "no_such_gazetteer.json")
