#!/usr/bin/env python3
"""Branded slide deck renderer (.pptx) behind python-pptx.

Turns the report view into a short, presentable deck: a title, a demand overview, a
winnable-clicks slide, a clusters table, and a how-to-read slide. Brand colours and
the brand name drive the look. The core-properties dates are pinned to the run
timestamp so that, together with the port's archive normalisation, two runs on the
same input produce a byte-identical file.

This path needs ``python-pptx`` (an optional backend). The caller checks
``ooxml_capabilities()`` and skips this format cleanly when the package is absent.
The rendering is exercised in continuous integration, where the backend installs;
locally, with no backend, the suite degrades to the HTML dashboard.
"""

from __future__ import annotations

from datetime import datetime
from pathlib import Path
from typing import Any

from .dashboard_html import _fmt_band, _fmt_int, _fmt_pct


def _rgb(color_hex: str):
    from pptx.dml.color import RGBColor

    return RGBColor.from_string(color_hex.lstrip("#"))


def _add_title(slide, text: str, brand: dict[str, Any], top_in: float = 0.4) -> None:
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(0.6), Inches(top_in), Inches(11.5), Inches(1.0))
    frame = box.text_frame
    frame.word_wrap = True
    run = frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = _rgb(brand["colors"]["primary"])


def _add_paragraph(frame, text: str, brand: dict[str, Any], size: int = 16, bold: bool = False, color: str | None = None) -> None:
    from pptx.util import Pt

    para = frame.add_paragraph()
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color or brand["colors"]["text"])


def _blank_slide(prs):
    # Layout 6 is the blank layout in the default template.
    return prs.slides.add_slide(prs.slide_layouts[6])


def _overview_slide(prs, view: dict[str, Any], brand: dict[str, Any]) -> None:
    from pptx.util import Inches

    slide = _blank_slide(prs)
    _add_title(slide, "Demand overview", brand)
    corpus = view["corpus"]
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(11.5), Inches(4.5))
    frame = box.text_frame
    frame.word_wrap = True
    lines = [
        (f"Keywords analysed: {_fmt_int(corpus['total_keywords'])}", True),
        (f"AI Overview eligible: {_fmt_pct(corpus['aio_eligibility_share'])}", False),
        (f"Generative-engine opportunity: {_fmt_pct(corpus['geo_opportunity_share'])}", False),
        (f"Demand opportunity score: {corpus['demand_opportunity_score']}", False),
    ]
    first = True
    for text, bold in lines:
        if first:
            run = frame.paragraphs[0].add_run()
            run.text = text
            from pptx.util import Pt

            run.font.size = Pt(18)
            run.font.bold = bold
            first = False
        else:
            _add_paragraph(frame, text, brand, size=18, bold=bold)
    _add_paragraph(frame, "Search intent split:", brand, size=15, bold=True, color=brand["colors"]["secondary"])
    for row in corpus["intent_split"]:
        _add_paragraph(frame, f"  {row['intent'].replace('_', ' ')}: {row['count']}", brand, size=14)


def _winnable_slide(prs, view: dict[str, Any], brand: dict[str, Any]) -> None:
    from pptx.util import Inches

    winnable = view.get("winnable")
    if not winnable:
        return
    slide = _blank_slide(prs)
    _add_title(slide, "Winnable clicks", brand)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(11.5), Inches(4.5))
    frame = box.text_frame
    frame.word_wrap = True
    run = frame.paragraphs[0].add_run()
    run.text = f"Portfolio winnable band per month: {_fmt_band(winnable.get('portfolio_winnable_band'))}"
    from pptx.util import Pt

    run.font.size = Pt(20)
    run.font.bold = True
    _add_paragraph(frame, f"Modelled current clicks: {_fmt_int(winnable.get('current_clicks_estimate'))}", brand, size=15)
    if "observed_current_clicks" in winnable:
        _add_paragraph(
            frame,
            f"Observed current clicks (Live Wire): {_fmt_int(winnable.get('observed_current_clicks'))}",
            brand, size=15, color=brand["colors"]["accent"],
        )
    _add_paragraph(frame, "By intent:", brand, size=15, bold=True, color=brand["colors"]["secondary"])
    for row in winnable["by_intent"]:
        _add_paragraph(
            frame,
            f"  {row['intent'].replace('_', ' ')}: {_fmt_band(row['winnable_band'])}",
            brand, size=14,
        )


def _clusters_slide(prs, view: dict[str, Any], brand: dict[str, Any]) -> None:
    from pptx.util import Inches, Pt

    slide = _blank_slide(prs)
    _add_title(slide, "Clusters by demand", brand)
    rows = view["top_clusters"][:10]
    headers = ["Cluster", "Volume", "Intent", "Readiness", "Winnable clicks"]
    table_shape = slide.shapes.add_table(
        len(rows) + 1, len(headers), Inches(0.6), Inches(1.6), Inches(12.0), Inches(0.4 * (len(rows) + 1))
    )
    table = table_shape.table
    for col, text in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(brand["colors"]["primary"])
        para = cell.text_frame.paragraphs[0]
        para.runs[0].font.size = Pt(12)
        para.runs[0].font.bold = True
        para.runs[0].font.color.rgb = _rgb("#ffffff")
    for r, c in enumerate(rows, start=1):
        values = [
            str(c["head"]),
            _fmt_int(c.get("volume_total")),
            str(c.get("dominant_intent", "")).replace("_", " "),
            str(c.get("expected_readiness", "n/a")),
            _fmt_band(c.get("winnable_band")),
        ]
        for col, value in enumerate(values):
            cell = table.cell(r, col)
            cell.text = value
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)


def _provenance_slide(prs, view: dict[str, Any], brand: dict[str, Any]) -> None:
    from pptx.util import Inches, Pt

    notes = view.get("provenance_notes") or []
    if not notes:
        return
    slide = _blank_slide(prs)
    _add_title(slide, "How to read these figures", brand)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(11.5), Inches(4.5))
    frame = box.text_frame
    frame.word_wrap = True
    first = True
    for note in notes:
        if first:
            run = frame.paragraphs[0].add_run()
            run.text = f"- {note}"
            run.font.size = Pt(14)
            first = False
        else:
            _add_paragraph(frame, f"- {note}", brand, size=14)


def render_deck(view: dict[str, Any], brand: dict[str, Any], out_path: Path, timestamp: datetime) -> None:
    """Render the branded deck to ``out_path``. Requires python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide.
    title_slide = _blank_slide(prs)
    header = view["header"]
    title_text = header["label"] or f"{brand['name']} keyword and demand audit"
    _add_title(title_slide, title_text, brand, top_in=2.6)
    sub = title_slide.shapes.add_textbox(Inches(0.6), Inches(3.7), Inches(11.5), Inches(1.0)).text_frame
    srun = sub.paragraphs[0].add_run()
    srun.text = f"{brand['name']} | {brand['tagline']}"
    srun.font.size = Pt(16)
    srun.font.color.rgb = _rgb(brand["colors"]["muted"])

    _overview_slide(prs, view, brand)
    _winnable_slide(prs, view, brand)
    _clusters_slide(prs, view, brand)
    _provenance_slide(prs, view, brand)

    core = prs.core_properties
    core.author = brand["author"]
    core.title = title_text
    core.created = timestamp
    core.modified = timestamp

    prs.save(str(out_path))
